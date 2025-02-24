import streamlit as st
import cohere
import PyPDF2
import pandas as pd
from docx import Document
from pptx import Presentation
from PIL import Image
import os
import io
import fitz  # PyMuPDF
from pyngrok import ngrok
import threading
from config import config
import easyocr
import numpy as np

COHERE_API_KEY = config["COHERE_API_KEY"]
co = cohere.Client(COHERE_API_KEY)

# File parsers
def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def read_pdf(file_path):
    content = ""
    # Using PyMuPDF (fitz) to extract text from PDFs
    with fitz.open(file_path) as doc:
        for page in doc:
            content += page.get_text("text")
    return content

def read_excel(file_path):
    data = pd.read_excel(file_path)
    return data.to_string()

def read_csv(file_path):
    data = pd.read_csv(file_path)
    return data.to_string()

def read_docx(file_path):
    doc = Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def read_pptx(file_path):
    presentation = Presentation(file_path)
    return "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, 'text')])

def read_image(file_path):
    # Use EasyOCR to extract text from images
    image = Image.open(file_path)
    image_np = np.array(image)
    reader = easyocr.Reader(['en'])
    results = reader.readtext(image_np, detail=0)
    return " ".join(results)

def read_file(file_path):
    if file_path.endswith('.txt'):
        return read_txt(file_path)
    elif file_path.endswith('.pdf'):
        return read_pdf(file_path)
    elif file_path.endswith(('.xls', '.xlsx')):
        return read_excel(file_path)
    elif file_path.endswith('.csv'):
        return read_csv(file_path)
    elif file_path.endswith('.docx'):
        return read_docx(file_path)
    elif file_path.endswith('.pptx'):
        return read_pptx(file_path)
    elif file_path.endswith(('.jpg', '.jpeg', '.png', '.tiff', '.tif')):
        return read_image(file_path)
    else:
        return "Unsupported file type."

def chatbot_response(query, document_text):
    prompt = f"""
    Document content:\n\n{document_text}\n\n
    Question: {query}
    Answer concisely and directly.
    """
    response = co.generate(
        model="command-xlarge-nightly",
        prompt=prompt,
        max_tokens=150,
        temperature=0.7
    )
    return response.generations[0].text.strip()

# Streamlit App
def run_streamlit():
    st.title("Document-Based Chatbot")
    st.write("Upload a document, ask a question, and get answers!")

    # File upload widget
    uploaded_file = st.file_uploader("Upload your document", 
                                     type=['txt', 'pdf', 'docx', 'xlsx', 'csv', 'pptx', 'jpg', 'jpeg', 'png', 'tiff', 'tif'])

    if uploaded_file:
        # Save the file temporarily
        file_path = os.path.join("temp", uploaded_file.name)
        os.makedirs("temp", exist_ok=True)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        # Process the file based on its extension
        document_text = read_file(file_path)
        st.success("Document processed successfully!")

        # Input for user query
        user_query = st.text_input("Enter your question about the document:")
        if user_query:
            response = chatbot_response(user_query, document_text)
            st.write("### Chatbot Response:")
            st.write(response)

# Ngrok Integration (if needed for remote access)
def start_ngrok():
    ngrok.set_auth_token(config["NGROK_AUTH_TOKEN"])
    public_url = ngrok.connect(8501)
    print(f"Streamlit app is live at: {public_url}")

if __name__ == "__main__":
    # Start Ngrok in a separate thread
    thread = threading.Thread(target=start_ngrok)
    thread.start()

    # Run the Streamlit app
    run_streamlit()
